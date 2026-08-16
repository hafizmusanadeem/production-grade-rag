from pathlib import Path

import logfire
from pptx import Presentation

from .base import BaseLoader, LoadedDocument


class PPTXLoader(BaseLoader):
    @logfire.instrument("Load PPTX {file_path=}", extract_args=("file_path",))
    def load(self, file_path: str | Path) -> list[LoadedDocument]:
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"PPTX not found: {path}")

        presentation = Presentation(str(path))
        docs: list[LoadedDocument] = []

        for slide_num, slide in enumerate(presentation.slides, start=1):
            texts: list[str] = []

            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue
                text = shape.text.strip()
                if text:
                    texts.append(text)

            if not texts:
                continue

            metadata = self._base_metadata(path, "pptx")
            metadata.update({"slide": slide_num})
            docs.append(
                LoadedDocument(
                    page_content="\n".join(texts),
                    metadata=metadata,
                )
            )

        if not docs:
            raise ValueError(f"No extractable text found in PPTX: {path}")

        self._record_load_metrics(
            path,
            docs,
            file_type="pptx",
            slide_count=len(presentation.slides),
            extracted_slides=len(docs),
        )
        return docs
